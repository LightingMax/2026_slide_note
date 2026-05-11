import shutil
import uuid
from datetime import datetime, timezone
from pathlib import Path

from pptx import Presentation

from app.core.config import get_settings
from app.models.deck import Deck, Slide, SlideAsset
from app.services.ppt_renderer import render_deck_snapshots


MEDIA_EXTENSIONS = {
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".gif": "image",
    ".webp": "image",
    ".bmp": "image",
    ".mp3": "audio",
    ".wav": "audio",
    ".m4a": "audio",
    ".aac": "audio",
    ".mp4": "video",
    ".mov": "video",
    ".webm": "video",
}


def parse_pptx(source_path: Path, original_filename: str) -> Deck:
    deck_id = uuid.uuid4().hex
    deck_media_dir = get_settings().media_dir / deck_id
    deck_media_dir.mkdir(parents=True, exist_ok=True)

    copied_path = get_settings().upload_dir / f"{deck_id}_{original_filename}"
    shutil.copyfile(source_path, copied_path)

    presentation = Presentation(str(source_path))
    slides = []

    for index, slide in enumerate(presentation.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_parts.append(shape.text.strip())

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        text = "\n".join(part for part in text_parts if part)
        title = _first_non_empty_line(text) or f"Slide {index}"
        slides.append(
            Slide(
                id=f"slide-{index}",
                index=index,
                title=title,
                text=text,
                notes=notes,
                assets=_extract_slide_media(slide, deck_id, deck_media_dir),
            )
        )

    deck = Deck(
        id=deck_id,
        filename=original_filename,
        created_at=datetime.now(timezone.utc).isoformat(),
        slides=slides,
    )
    return render_deck_snapshots(deck, copied_path)


def _extract_slide_media(slide, deck_id: str, deck_media_dir: Path) -> list[SlideAsset]:
    assets: list[SlideAsset] = []
    seen: set[str] = set()
    for relationship in slide.part.rels.values():
        try:
            target_part = relationship.target_part
        except ValueError:
            continue

        path = Path(str(target_part.partname))
        extension = path.suffix.lower()
        if extension not in MEDIA_EXTENSIONS or str(path) in seen:
            continue

        seen.add(str(path))
        filename = f"{uuid.uuid4().hex}_{path.name}"
        target = deck_media_dir / filename
        target.write_bytes(target_part.blob)
        assets.append(
            SlideAsset(
                name=path.name,
                kind=MEDIA_EXTENSIONS[extension],
                url=f"/media/{deck_id}/{filename}",
                content_type=getattr(target_part, "content_type", None),
            )
        )
    return assets


def _first_non_empty_line(text: str) -> str:
    for line in text.splitlines():
        value = line.strip()
        if value:
            return value[:80]
    return ""
