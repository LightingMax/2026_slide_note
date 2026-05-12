import shutil
from pathlib import Path

from pptx import Presentation

from app.core.config import get_settings
from app.models.deck import Deck


def find_uploaded_pptx(deck_id: str) -> Path | None:
    for candidate in get_settings().upload_dir.glob(f"{deck_id}_*.pptx"):
        return candidate
    return None


def export_deck_with_notes(deck: Deck, source_path: Path) -> Path:
    export_dir = get_settings().data_dir / "exports"
    export_dir.mkdir(parents=True, exist_ok=True)
    target = export_dir / f"{deck.id}_slide-note.pptx"
    shutil.copyfile(source_path, target)

    presentation = Presentation(str(target))
    notes_by_index = {slide.index: slide.notes for slide in deck.slides}
    for index, slide in enumerate(presentation.slides, start=1):
        notes = notes_by_index.get(index, "")
        notes_slide = slide.notes_slide
        if notes_slide.notes_text_frame is not None:
            notes_slide.notes_text_frame.text = notes

    presentation.save(str(target))
    return target

